#!/usr/bin/env python3
"""Render a stakeholder ship brief as a .pptx in the Kinoped Keynote style.

Invoked by the `brief` skill (see SKILL.md next to this file) for the default
pptx format; `/brief md` renders Markdown instead and does not use this script.

Usage:
    python3 generate_brief_pptx.py spec.json -o docs/briefs/2026-07-27-brief.pptx

Requires python-pptx (pip install python-pptx). The design system is extracted
from Kinoped_Whats_New.key (13.333x7.5in, 16:9):
  bg light F5F8F7 / dark 0B1220; deep teal 0F766E, bright teal 5EEAD4;
  Georgia headings #16241F, Calibri body #3A4A45; rounded teal number badges;
  white cards (FFFFFF, 1pt E4ECEA border, ~0.045 corner radius).

Spec shape (JSON):
{
  "brand": "KINOPED",                  # small-caps brand line on the title slide
  "product": "Platform",
  "title": "What's New",
  "subtitle": "Foundation Platform - Release Overview",
  "date_range": "17 Jun - 27 Jul 2026",
  "strip": "patterns - sessions - edge sim",   # optional one-line feature strip
  "features": [
    {
      "area": "SESSIONS",              # EYEBROW, uppercased
      "title": "Gait metrics on session rows",
      "lead": "One or two sentences on what shipped and why it matters.",
      "cards": [                       # exactly 3 reads best; 2-3 supported
        {"n": "1", "title": "What", "body": "..."},
        {"n": "2", "title": "How", "body": "..."},
        {"n": "3", "title": "Impact", "body": "..."}
      ],
      "caption": "Session list, right sidebar",  # optional '>' demo caption
      "image": "path/to/screenshot.png"          # optional; replaces card row
    }
  ],
  "next": [                            # optional dark "looking ahead" slide
    {"title": "...", "desc": "..."}
  ]
}
"""
import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls, qn
except ImportError:
    sys.exit("python-pptx is required: pip install python-pptx (use a venv)")

# ---- design tokens (extracted from the Kinoped What's New deck) ------------
BG_LIGHT = "F5F8F7"
BG_DARK = "0B1220"
TEAL_DEEP = "0F766E"
TEAL_BRIGHT = "5EEAD4"
INK = "16241F"
BODY = "3A4A45"
CARD_BORDER = "E4ECEA"
LIGHT_ON_DARK = "E6EEEC"
MUTE_ON_DARK = "9DB0AC"
GEORGIA = "Georgia"
CALIBRI = "Calibri"

SLIDE_W = 13.333
SLIDE_H = 7.5


def set_bg(slide, hexcolor):
    cSld = slide._element.find(qn("p:cSld"))
    bg = parse_xml(
        f'<p:bg {nsdecls("p", "a")}><p:bgPr>'
        f'<a:solidFill><a:srgbClr val="{hexcolor}"/></a:solidFill>'
        f"<a:effectLst/></p:bgPr></p:bg>"
    )
    cSld.insert(0, bg)


def textbox(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            wrap=True, space_after=None, line_spacing=None):
    """runs: list of paragraphs; each paragraph is a list of
    (text, font, size, bold, color) run tuples."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after is not None:
            p.space_after = Pt(space_after)
        if line_spacing is not None:
            p.line_spacing = line_spacing
        for (text, font, size, bold, color) in para:
            r = p.add_run()
            r.text = text
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = RGBColor.from_string(color)
    return tb


def rounded(slide, l, t, w, h, fill, line=None, line_w=None, radius=0.05):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(line_w or 1.0)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def oval(slide, l, t, w, h, fill):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def badge(slide, label, l=0.55, t=0.50, size=0.62, font_size=18):
    sp = rounded(slide, l, t, size, size, TEAL_DEEP, radius=0.16129)
    tf = sp.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = GEORGIA
    r.font.size = Pt(font_size)
    r.font.bold = True
    r.font.color.rgb = RGBColor.from_string("FFFFFF")
    return sp


def header(slide, num_label, eyebrow, title):
    badge(slide, num_label)
    textbox(slide, 1.32, 0.50, 11.6, 0.32,
            [[(eyebrow.upper(), CALIBRI, 11.5, True, TEAL_DEEP)]])
    textbox(slide, 1.32, 0.86, 11.6, 0.80,
            [[(title, GEORGIA, 30, True, INK)]])


def dim_ring(slide, l=9.7, t=4.4, d=6.4, alpha=14000):
    ring = oval(slide, l, t, d, d, TEAL_DEEP)
    sf = ring.fill.fore_color._xFill.find(qn("a:srgbClr"))
    if sf is not None:
        sf.append(parse_xml(f'<a:alpha {nsdecls("a")} val="{alpha}"/>'))
    return ring


# ---- slides ----------------------------------------------------------------
def title_slide(prs, layout, spec):
    s = prs.slides.add_slide(layout)
    set_bg(s, BG_DARK)
    dim_ring(s, 8.9, 3.6, 7.4)
    textbox(s, 0.9, 1.55, 11.0, 0.4,
            [[(spec.get("brand", "").upper(), CALIBRI, 15, True, TEAL_BRIGHT)]])
    textbox(s, 0.87, 2.05, 11.6, 2.0,
            [[(spec.get("product", ""), GEORGIA, 40, True, "FFFFFF")],
             [(spec.get("title", "What's New"), GEORGIA, 40, True, TEAL_BRIGHT)]],
            line_spacing=1.05)
    sub = spec.get("subtitle", "")
    rng = spec.get("date_range", "")
    subline = " · ".join(x for x in (sub, rng) if x)
    textbox(s, 0.9, 4.35, 11.6, 0.5,
            [[(subline, CALIBRI, 15, False, MUTE_ON_DARK)]])
    strip = spec.get("strip", "")
    if strip:
        textbox(s, 0.9, 6.55, 12.0, 0.4,
                [[(strip, CALIBRI, 11.5, False, MUTE_ON_DARK)]])
    return s


def glance_slide(prs, layout, spec):
    """'00 THIS RELEASE / What's new, at a glance' - one row per feature."""
    s = prs.slides.add_slide(layout)
    set_bg(s, BG_LIGHT)
    header(s, "00", "This release", "What's new, at a glance")
    feats = spec.get("features", [])
    # single column up to 5, two columns beyond
    if len(feats) <= 5:
        cols = [(0.9, feats)]
        col_w = 11.3
    else:
        mid = (len(feats) + 1) // 2
        cols = [(0.9, feats[:mid]), (6.95, feats[mid:])]
        col_w = 5.4
    top, step = 2.05, 0.98
    for x, chunk in cols:
        for i, f in enumerate(chunk):
            y = top + i * step
            oval(s, x, y + 0.09, 0.16, 0.16, TEAL_DEEP)
            textbox(s, x + 0.38, y, col_w, 0.36,
                    [[(f["title"], CALIBRI, 15, True, INK)]])
            one = f.get("one_liner") or f.get("lead", "")
            textbox(s, x + 0.38, y + 0.36, col_w, 0.5,
                    [[(one[:160], CALIBRI, 11, False, BODY)]], line_spacing=1.05)
    return s


def feature_slide(prs, layout, num_label, f, repo_root):
    s = prs.slides.add_slide(layout)
    set_bg(s, BG_LIGHT)
    header(s, num_label, f.get("area", ""), f["title"])
    textbox(s, 0.57, 1.74, 12.2, 0.62,
            [[(f.get("lead", ""), CALIBRI, 14, False, BODY)]], line_spacing=1.05)

    image = f.get("image")
    img_path = (repo_root / image) if image else None
    if img_path and img_path.exists():
        # screenshot panel: white card frame with the image inside, caption below
        rounded(s, 0.55, 2.55, 12.23, 4.15, "FFFFFF", line=CARD_BORDER, line_w=1.0,
                radius=0.03)
        pic = s.shapes.add_picture(str(img_path), Inches(0.8), Inches(2.8))
        # scale to fit the panel
        max_w, max_h = Inches(11.7), Inches(3.6)
        ratio = min(max_w / pic.width, max_h / pic.height, 1.0)
        pic.width = int(pic.width * ratio)
        pic.height = int(pic.height * ratio)
        pic.left = Inches(0.55) + int((Inches(12.23) - pic.width) / 2)
        pic.top = Inches(2.65) + int((Inches(3.9) - pic.height) / 2)
    else:
        cards = f.get("cards", [])[:3]
        xs = [0.55, 4.71, 8.87][:len(cards)]
        cw, ct, ch = 3.91, 2.62, 3.86
        for x, c in zip(xs, cards):
            rounded(s, x, ct, cw, ch, "FFFFFF", line=CARD_BORDER, line_w=1.0,
                    radius=0.045)
            ic = rounded(s, x + 0.24, ct + 0.26, 0.58, 0.58, TEAL_DEEP, radius=0.22)
            itf = ic.text_frame
            itf.vertical_anchor = MSO_ANCHOR.MIDDLE
            for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
                setattr(itf, m, 0)
            ip = itf.paragraphs[0]
            ip.alignment = PP_ALIGN.CENTER
            ir = ip.add_run()
            ir.text = str(c.get("n", ""))
            ir.font.name = GEORGIA
            ir.font.size = Pt(17)
            ir.font.bold = True
            ir.font.color.rgb = RGBColor.from_string("FFFFFF")
            textbox(s, x + 0.24, ct + 1.04, cw - 0.48, 0.6,
                    [[(c.get("title", ""), GEORGIA, 14, True, INK)]], line_spacing=1.0)
            textbox(s, x + 0.24, ct + 1.62, cw - 0.48, ch - 1.8,
                    [[(c.get("body", ""), CALIBRI, 11.5, False, BODY)]],
                    line_spacing=1.12)
    caption = f.get("caption")
    if caption:
        textbox(s, 0.57, 6.78, 12.2, 0.4,
                [[("▶  " + caption, CALIBRI, 11.5, True, TEAL_DEEP)]])
    return s


def next_slide(prs, layout, rows):
    s = prs.slides.add_slide(layout)
    set_bg(s, BG_DARK)
    dim_ring(s)
    textbox(s, 0.85, 0.66, 11.0, 0.32,
            [[("LOOKING AHEAD · PROPOSED", CALIBRI, 12, True, TEAL_BRIGHT)]])
    textbox(s, 0.83, 1.04, 11.6, 0.9,
            [[("What we could build next", GEORGIA, 33, True, "FFFFFF")]])
    textbox(s, 0.86, 1.96, 11.6, 0.5,
            [[("Candidate initiatives for the next phase — not yet shipped.",
               CALIBRI, 14, False, MUTE_ON_DARK)]])
    top, step = 2.78, 1.04
    for i, row in enumerate(rows[:4]):
        y = top + i * step
        oval(s, 0.9, y + 0.10, 0.18, 0.18, TEAL_BRIGHT)
        textbox(s, 1.32, y, 11.3, 0.4,
                [[(row["title"], CALIBRI, 16, True, LIGHT_ON_DARK)]])
        textbox(s, 1.32, y + 0.40, 11.3, 0.4,
                [[(row.get("desc", ""), CALIBRI, 11.5, False, MUTE_ON_DARK)]])
    return s


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("spec", help="path to the JSON content spec")
    ap.add_argument("-o", "--out", required=True, help="output .pptx path")
    args = ap.parse_args()

    spec = json.loads(Path(args.spec).read_text())
    out = Path(args.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    repo_root = Path.cwd()

    prs = Presentation()
    prs.slide_width = Emu(int(Inches(SLIDE_W)))
    prs.slide_height = Emu(int(Inches(SLIDE_H)))
    layout = prs.slide_layouts[6]  # blank

    title_slide(prs, layout, spec)
    glance_slide(prs, layout, spec)
    for i, f in enumerate(spec.get("features", []), start=1):
        feature_slide(prs, layout, f"{i:02d}", f, repo_root)
    if spec.get("next"):
        next_slide(prs, layout, spec["next"])

    prs.save(out)
    print(f"Saved {out} ({2 + len(spec.get('features', [])) + (1 if spec.get('next') else 0)} slides)")


if __name__ == "__main__":
    main()
